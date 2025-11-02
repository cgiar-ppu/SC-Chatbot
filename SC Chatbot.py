import os
import re
from dataclasses import dataclass
from typing import List, Dict, Optional, Tuple

import streamlit as st
from sklearn.metrics.pairwise import cosine_similarity

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openai import OpenAI

# --- CGIAR Theme --------------------------------------------------------------
import html  # to escape text in chips

import pandas as pd
import pickle
import hashlib
from dataclasses import asdict
import numpy as np
import torch
from sentence_transformers import SentenceTransformer

from datetime import datetime
import json
from urllib.parse import urlparse

device = 'cuda' if torch.cuda.is_available() else 'cpu'

CGIAR_COLORS = {
    "primary": "#1F5A48",
    "hover": "#285f49",
    "link": "#0B66C3",
    "sidebar": "#FBF2D2",
    "bg": "#FFFFFF",
    "panel": "#FFFFFF",
    "text": "#1F5A48",
    "muted": "#4A5568",
    "border": "#1F5A48",  # Updated
    "callout": "#f2faf6",
    "callout-border": "#1F5A48",
    "separator": "#1F5A48",
}

def apply_cgiar_theme():
    st.markdown(f"""
<style>
    /* Typography */
    @import url('https://fonts.googleapis.com/css2?family=Montserrat:wght@400;500;600;700&family=Open+Sans:wght@400;600&display=swap');

    :root {{
        --brand-primary: {CGIAR_COLORS["primary"]};
        --brand-primary-strong: {CGIAR_COLORS["hover"]};
        --brand-blue: {CGIAR_COLORS["link"]};
        --sidebar-bg: {CGIAR_COLORS["sidebar"]};

        --bg: {CGIAR_COLORS["bg"]};
        --panel: {CGIAR_COLORS["panel"]};
        --text: {CGIAR_COLORS["text"]};
        --muted: {CGIAR_COLORS["muted"]};
        --border: {CGIAR_COLORS["border"]};
        --callout: {CGIAR_COLORS["callout"]};
        --callout-border: {CGIAR_COLORS["callout-border"]};
        --separator: {CGIAR_COLORS["separator"]};
    }}

    /* Ensure body font size */
    .stApp {{
        font-family: 'Open Sans', sans-serif;
        background: var(--bg);
        color: var(--text);
        font-size: 16px; /* >=16px */
    }}

    /* Titles */
    h1, h2, h3, h4, h5, h6 {{
        font-family: 'Montserrat', sans-serif;
    }}
    .brand-hero h1 {{
        font-size: 1.1875rem;
    }}

    /* Sidebar - hide completely */
    section[data-testid="stSidebar"] {{
        display: none !important;
    }}

    /* Chat card with green border and text */
    .chat-card {{
        background: #FFFFFF;
        border: 1px solid #1F5A48;
        border-radius: 12px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.04);
        padding: 1.25rem;
        margin: 1rem 0;
        color: #1F5A48;
    }}

    /* Ensure search bar is green */
    .stTextInput > div > div > input {{
        background-color: #1F5A48 !important;
        color: #FFFFFF !important;
        border: 1px solid #1F5A48 !important;
        padding: 0.5rem 1rem;
    }}

    .stTextInput > div > div > input::placeholder {{
        color: rgba(255,255,255,0.7);
    }}

    .stTextInput > div > div > input:focus {{
        border-color: #285f49 !important;
        box-shadow: 0 0 0 2px rgba(40,95,73,0.2);
    }}

    /* Streamlit header bar */
    [data-testid="stHeader"] {{
        background-color: #FFFFFF !important;
    }}

    /* Search input hover */
    .stTextInput > div > div > input:hover {{
        background-color: #285f49 !important; /* Darker green on hover */
        border-color: #285f49 !important;
        box-shadow: 0 0 0 2px rgba(40,95,73,0.2);
    }}

    /* Override Streamlit button defaults completely */
    button[kind="primary"] {{
        background-color: #FFFFFF !important;
        color: #1F5A48 !important;
        border: 1px solid #1F5A48 !important;
        transition: none !important;
    }}

    button[kind="primary"]:hover {{
        background-color: #FFFFFF !important;
        color: #1F5A48 !important;
        border: 1px solid #1F5A48 !important;
        box-shadow: none !important;
    }}

    /* Green callout with full border */
    .brand-hero {{
        background: #f2faf6;
        border: 1px solid #1F5A48;
        border-radius: 12px;
        padding: 1.25rem;
        margin-bottom: 1rem;
        color: #1F5A48;
    }}

    /* Top Sources table */
    .leaderboard thead th {{
        font-family: 'Montserrat', sans-serif;
        border-bottom: 2px solid #e8ece8;
    }}

    .leaderboard tbody tr {{
        border-top: 1px solid #e8ece8;
    }}

    /* Accessibility */
    :focus-visible {{
        outline: 2px solid #0B66C3;
    }}

    /* Custom spinner styling */
    .custom-spinner {{
        width: 50px;
        height: 50px;
        border: 8px solid rgba(0, 0, 0, 0.1);
        border-top-color: #1F5A48;
        border-radius: 50%;
        animation: spin 1s linear infinite;
    }}
    @keyframes spin {{
        to {{ transform: rotate(360deg); }}
    }}

    /* Main container */
    .main .block-container {{
        max-width: 1100px;
        padding-top: 1.25rem;
    }}

    /* Slider */
    .stSlider [data-baseweb="slider"] > div:first-child {{
        color: var(--brand-primary);
    }}

    /* Subtle alerts */
    .stAlert {{
        border-left: 4px solid var(--brand-accent);
    }}

    /* Compact metrics */
    .metric-row > div > div {{
        background: var(--panel);
        border: 1px solid var(--border);
        border-radius: 12px;
        padding: .75rem;
    }}

    /* Footer */
    .app-footer {{
        text-align: center;
        color: var(--muted);
        font-size: .9rem;
        margin: 1rem 0 2rem;
    }}
</style>
""", unsafe_allow_html=True)

@st.cache_resource
def get_embedding_model() -> SentenceTransformer:
    model_name = 'all-MiniLM-L6-v2'
    model_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'models', 'sentence_transformer')
    os.makedirs(model_dir, exist_ok=True)
    try:
        model = SentenceTransformer(model_dir)
    except Exception:
        model = SentenceTransformer(model_name)
        try:
            model.save(model_dir)
        except Exception:
            pass
    return model.to(device)

@dataclass
class Chunk:
    text: str
    source_path: str
    source_name: str
    kind: str  # pdf | docx | pptx
    location: str  # e.g., "page 3", "slide 2", "section X / paragraph 12"
    id: str = ""


def normalize_whitespace(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def split_sentences(text: str) -> List[str]:
    if not text:
        return []
    parts = re.split(r"(?<=[\.\!\?])\s+", text)
    sentences: List[str] = []
    for p in parts:
        p = p.strip()
        if len(p) > 0:
            sentences.append(p)
    if len(sentences) <= 1:
        lines = [ln.strip() for ln in re.split(r"[\n\r]+", text) if ln.strip()]
        if len(lines) > len(sentences):
            sentences = lines
    return sentences


def split_text_into_overlapping_parts(text: str, num_parts: int = 4, overlap_fraction: float = 0.25) -> List[str]:
    words = text.split()
    total_words = len(words)
    if total_words == 0:
        return []
    denominator = num_parts - (num_parts - 1) * overlap_fraction if num_parts > 1 else 1
    part_words = int(total_words / denominator) if denominator > 0 else total_words
    overlap_words = int(part_words * overlap_fraction)
    step = part_words - overlap_words
    parts = []
    start = 0
    for i in range(num_parts):
        end = min(start + part_words, total_words)
        part = ' '.join(words[start:end])
        if part:
            parts.append(part)
        if end >= total_words:
            break
        start += step
    return parts


def read_pdf_chunks(path: str) -> List[Chunk]:
    chunks: List[Chunk] = []
    try:
        reader = PdfReader(path)
        for index, page in enumerate(reader.pages, start=1):
            page_text = normalize_whitespace(page.extract_text() or "")
            if not page_text:
                continue
            parts = split_text_into_overlapping_parts(page_text, 4, 0.25)
            for part_idx, part in enumerate(parts, start=1):
                chunks.append(
                    Chunk(
                        text=part,
                        source_path=path,
                        source_name=os.path.basename(path),
                        kind="pdf",
                        location=f"page {index} part {part_idx}",
                        id=f"pdf:{os.path.basename(path)}:p{index}pt{part_idx}",
                    )
                )
        return chunks
    except Exception:
        return []


def read_docx_chunks(path: str) -> List[Chunk]:
    try:
        doc = DocxDocument(path)
        unit_texts: List[str] = []
        unit_locations: List[str] = []
        current_section: Optional[str] = None
        for paragraph_index, paragraph in enumerate(doc.paragraphs, start=1):
            text = normalize_whitespace(paragraph.text)
            if not text:
                continue
            style_name = getattr(paragraph.style, "name", "") or ""
            if style_name.lower().startswith("heading") or style_name.lower().startswith("título"):
                current_section = text
            location = f"section '{current_section}' paragraph {paragraph_index}" if current_section else f"paragraph {paragraph_index}"
            unit_texts.append(text)
            unit_locations.append(location)
        for t_idx, table in enumerate(doc.tables, start=1):
            for r_idx, row in enumerate(table.rows, start=1):
                cells = [normalize_whitespace(cell.text or "") for cell in row.cells]
                cells = [c for c in cells if c]  # quita celdas vacías
                if cells:
                    unit_texts.append(" | ".join(cells))
                    unit_locations.append(
                        f"{'section ' + repr(current_section) + ' ' if current_section else ''}table {t_idx} row {r_idx}"
                    )
        return create_overlapped_chunks(unit_texts, unit_locations, os.path.basename(path), "docx", path)
    except Exception:
        pass
    return []


def read_pptx_chunks(path: str) -> List[Chunk]:
    chunks: List[Chunk] = []
    try:
        prs = Presentation(path)
        for slide_index, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        txt = "\n".join(p.text for p in shape.text_frame.paragraphs)
                        txt = normalize_whitespace(txt)
                        if txt:
                            texts.append(txt)
                except Exception:
                    continue
            slide_text = normalize_whitespace("\n".join(texts))
            if not slide_text:
                continue
            parts = split_text_into_overlapping_parts(slide_text, 4, 0.25)
            for part_idx, part in enumerate(parts, start=1):
                chunks.append(
                    Chunk(
                        text=part,
                        source_path=path,
                        source_name=os.path.basename(path),
                        kind="pptx",
                        location=f"slide {slide_index} part {part_idx}",
                        id=f"pptx:{os.path.basename(path)}:s{slide_index}pt{part_idx}",
                    )
                )
        return chunks
    except Exception:
        return []


def read_excel_chunks(path: str) -> List[Chunk]:
    try:
        df = pd.read_excel(path, engine='openpyxl')
        required_columns = ['Decision', 'Date(s)']
        if not all(col in df.columns for col in required_columns):
            return []
        chunks = []
        basename = os.path.basename(path)
        for index, row in df.iterrows():
            decision = str(row['Decision']).strip()
            dates = str(row['Date(s)']).strip()
            text = f"Decision: {decision}\nDate(s): {dates}"
            location = f"row {index + 1}"
            chunk_id = f"xlsx:{basename}:row{index+1}"
            chunks.append(
                Chunk(
                    text=text,
                    source_path=path,
                    source_name=basename,
                    kind="xlsx",
                    location=location,
                    id=chunk_id,
                )
            )
        return chunks
    except Exception:
        return []


def create_overlapped_chunks(unit_texts: List[str], unit_locations: List[str], basename: str, kind: str, path: str, window_size: int = 4, overlap: int = 1) -> List[Chunk]:
    chunks: List[Chunk] = []
    step = window_size - overlap
    for i in range(0, len(unit_texts), step):
        end = i + window_size
        slice_texts = unit_texts[i:end]
        slice_locations = unit_locations[i:end] if end <= len(unit_locations) else unit_locations[i:]
        if len(slice_texts) == 0:
            continue
        merge_len = len(slice_texts)
        if merge_len < 3 and chunks:
            # merge to last
            last = chunks[-1]
            last.text += ' ' + ' '.join(slice_texts)
            last_loc_end = slice_locations[-1] if slice_locations else ""
            last.location = last.location.rsplit(" to ", 1)[0] + f" to {last_loc_end}"
            # update id
            parts = last.id.split(':')
            if len(parts) == 3 and parts[2].startswith('u'):
                range_part = parts[2][1:]
                start_str, end_str = range_part.split('-')
                start = int(start_str)
                old_end = int(end_str)
                new_end = old_end + merge_len
                last.id = f"{parts[0]}:{parts[1]}:u{start}-{new_end}"
            continue
        chunk_text = ' '.join(slice_texts)
        loc_start = slice_locations[0]
        loc_end = slice_locations[-1]
        chunk_loc = f"{loc_start} to {loc_end}"
        start_idx = i + 1
        end_idx = i + merge_len
        chunk_id = f"{kind}:{basename}:u{start_idx}-{end_idx}"
        chunks.append(
            Chunk(
                text=chunk_text,
                source_path=path,
                source_name=basename,
                kind=kind,
                location=chunk_loc,
                id=chunk_id,
            )
        )
    return chunks


def load_corpus(root_dir: str) -> List[Chunk]:
    supported_ext = {".pdf", ".docx", ".pptx", ".xlsx"}
    chunks: List[Chunk] = []
    latest_updates_dir = os.path.join(root_dir, 'Latest Updates')
    for dirpath, _, filenames in os.walk(latest_updates_dir):
        for fname in filenames:
            ext = os.path.splitext(fname)[1].lower()
            if ext not in supported_ext:
                continue
            abspath = os.path.join(dirpath, fname)
            try:
                if os.path.getsize(abspath) == 0:
                    continue
            except Exception:
                continue
            new_chunks = []
            try:
                if ext == ".pdf":
                    new_chunks = read_pdf_chunks(abspath)
                elif ext == ".docx":
                    new_chunks = read_docx_chunks(abspath)
                elif ext == ".pptx":
                    new_chunks = read_pptx_chunks(abspath)
                elif ext == ".xlsx":
                    new_chunks = read_excel_chunks(abspath)
            except Exception:
                continue
            chunks.extend(new_chunks)
    return chunks


def build_embeddings_index(chunks: List[Chunk]) -> np.ndarray:
    model = get_embedding_model()
    texts = [c.text for c in chunks]
    if len(texts) == 0:
        try:
            dim = model.get_sentence_embedding_dimension()
        except Exception:
            dim = 384
        return np.empty((0, dim), dtype=np.float32)
    with st.spinner("Building semantic embeddings..."):
        embs = model.encode(texts, show_progress_bar=True, device=device, normalize_embeddings=False)
    return np.asarray(embs, dtype=np.float32)

def rank_chunks(query: str, model: SentenceTransformer, embeddings: np.ndarray, chunks: List[Chunk], top_k: int = 25) -> List[Tuple[Chunk, float]]:
    if not query.strip() or embeddings.shape[0] == 0:
        return []
    q_vec = model.encode([query], device=device, normalize_embeddings=False)
    sims = cosine_similarity(q_vec, embeddings)[0]

    scored = [(idx, float(s)) for idx, s in enumerate(sims) if s > 0.0]

    scored.sort(key=lambda x: x[1], reverse=True)
    out: List[Tuple[Chunk, float]] = []
    for idx, s in scored[: max(top_k * 2, top_k)]:
        out.append((chunks[idx], s))
        if len(out) >= top_k:
            break
    return out


def extract_relevant_sentences(query: str, texts: List[str], max_sentences: int = 6) -> List[str]:
    query_terms = [t for t in re.split(r"\W+", query.lower()) if len(t) > 2]
    candidates: List[Tuple[str, float]] = []
    for block in texts:
        for sent in split_sentences(block):
            low = sent.lower()
            if not low:
                continue
            tf = sum(low.count(t) for t in query_terms)
            if tf == 0:
                continue
            length_penalty = 1.0 + max(0, (len(sent) - 300) / 300.0)
            score = tf / length_penalty
            candidates.append((sent.strip(), score))
    candidates.sort(key=lambda x: x[1], reverse=True)
    unique: List[str] = []
    seen = set()
    for sent, _ in candidates:
        key = sent[:120]
        if key in seen:
            continue
        seen.add(key)
        unique.append(sent)
        if len(unique) >= max_sentences:
            break
    return unique


def load_name_to_link(project_root: str) -> Dict[str, str]:
    excel_path = os.path.join(project_root, 'Docs & Links', 'Docs & Links - Chatbot P&R Hub.xlsx')
    if not os.path.exists(excel_path):
        return {}
    df = pd.read_excel(excel_path)
    name_col = 'NAME DOCUMENT'
    link_col = 'LINKS'
    if name_col not in df.columns or link_col not in df.columns:
        return {}
    mapping = {}
    for _, row in df.iterrows():
        name = str(row[name_col]).strip()
        link = str(row[link_col]).strip()
        if name and link:
            mapping[name] = link
    return mapping


def compose_answer(query: str, ranked: List[Tuple[Chunk, float]], name_to_link: Dict[str, str] = {}) -> Tuple[str, List[str]]:
    if not ranked:
        msg = (
            "Not found in the available information. "
            "A specific reference (document/page or section) would be needed. "
            "Verify the document name or try other keywords."
        )
        return msg, []
    texts = [c.text for c, _ in ranked]
    sentences = extract_relevant_sentences(query, texts, max_sentences=6)

    if len(sentences) == 0 and texts:
        sentences = split_sentences(texts[0])[:6]

    if len(sentences) == 0:
        msg = (
            "Not found in the available information. "
            "There are no relevant snippets for the current query."
        )
        return msg, []
    if len(sentences) < 3 and len(ranked) >= 2:
        extra_sentences = []
        for c, _ in ranked:
            sents = split_sentences(c.text)
            for s in sents:
                if len(s.strip()) > 40 and s.strip() not in sentences:
                    extra_sentences.append(s.strip())
                if len(sentences) + len(extra_sentences) >= 3:
                    break
            if len(sentences) + len(extra_sentences) >= 3:
                break
        sentences.extend(extra_sentences[: max(0, 3 - len(sentences))])
    sentences = sentences[:6]
    answer = " ".join(sentences).strip()
    return answer, [f"{name_to_link.get(c.source_name, c.source_name)} — {c.location} — {c.id}" for c, _ in ranked]


def format_sources_lines(ranked: List[Tuple[Chunk, float]] , max_items: int = 10, name_to_link: Dict[str, str] = {}) -> List[str]:
    lines: List[str] = []
    seen = set()
    for c, _ in ranked:
        entry = f"{name_to_link.get(c.source_name, c.source_name)} — {c.location} — {c.id}"
        if entry in seen:
            continue
        seen.add(entry)
        lines.append(entry)
        if len(lines) >= max_items:
            break
    return lines


def call_openai_generate(query: str, ranked: List[Tuple[Chunk, float]], max_sentences: int = 5, custom_system_msg: Optional[str] = None, name_to_link: Dict[str, str] = {}) -> Tuple[Optional[str], Optional[str]]:

    selected = ranked
    context_blocks: List[str] = []
    for c, _ in selected:
        link = name_to_link.get(c.source_name, '')
        context_blocks.append(
            f"[ID: {c.id}]\nLink: {link if link else 'N/A'}\nFile: {c.source_name}\nLocation: {c.location}\nContent: {c.text}"
        )
    context = "\n\n---\n\n" + "\n\n---\n\n".join(context_blocks) if context_blocks else ""

    base_msg = custom_system_msg if custom_system_msg is not None else "You are a helpful assistant that answer user questions in a detailed and comprehensive way."

    system_msg = base_msg + "\n\nAlways base your answer strictly on the provided context.\nProvide a detailed and comprehensive answer to the question.\nAfter the answer, add a section 'Used Sources:' followed by a comma-separated list of the IDs from the chunks you used in formulating the answer (e.g., pdf:doc1.pdf:p1pt1, docx:doc2.docx:u1-4). Only list chunks that directly contributed. Do not make up IDs."

    user_msg = (
        f"Question: {query}\n\nContext:{context}"
    )

    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": user_msg},
    ]
    full_input = json.dumps(messages, ensure_ascii=False)

    try:
        api_key = os.environ.get("OPENAI_API_KEY", "").strip()
        if not api_key:
            return (None, full_input)
        client = OpenAI(api_key=api_key)
        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=messages,
                temperature=0,
                max_tokens=5000,
            )
            return ((resp.choices[0].message.content or "").strip(), full_input)
        except Exception:
            try:
                resp2 = client.responses.create(
                    model="gpt-4o-mini",
                    input=messages,
                    temperature=0,
                    max_output_tokens=5000,
                )
                if hasattr(resp2, "output") and resp2.output and hasattr(resp2.output[0], "content"):
                    parts = resp2.output[0].content
                    if parts and hasattr(parts[0], "text"):
                        return ((parts[0].text or "").strip(), full_input)
            except Exception:
                return (None, full_input)
    except Exception:
        return (None, full_input)

    return (None, full_input)


def parse_used_sources(ai_answer: str) -> Tuple[str, List[str]]:
    if 'Used Sources:' not in ai_answer:
        return ai_answer.strip(), []
    parts = ai_answer.split('Used Sources:', 1)
    answer_part = parts[0].strip()
    sources_str = parts[1].strip()
    used_ids = [s.strip() for s in sources_str.split(',') if s.strip()]
    return answer_part, used_ids


def dedupe_preserve_order(items: List[str], limit: int = 5) -> List[str]:
    out: List[str] = []
    seen = set()
    for it in items:
        if it in seen:
            continue
        seen.add(it)
        out.append(it)
        if len(out) >= limit:
            break
    return out

def make_clickable_entry(entry: str) -> str:
    parts = entry.split(" — ", 2)
    if len(parts) < 3:
        return html.escape(entry)
    link, location, cited_id = parts
    if link.startswith("http://") or link.startswith("https://"):
        escaped_link = html.escape(link)
        return f'<a href="{escaped_link}" target="_blank" rel="noopener noreferrer">{escaped_link}</a> — {html.escape(location)} — {html.escape(cited_id)}'
    else:
        return html.escape(entry)

def render_sources_pills(lines: List[str]):
    if not lines:
        st.markdown("<div class='sources-wrap'><span class='source-chip'>not specified</span></div>", unsafe_allow_html=True)
        return
    clickable_lines = [make_clickable_entry(line) for line in lines]
    pills = "".join(f"<span class='source-chip'>{pill}</span>" for pill in clickable_lines)
    st.markdown(f"<div class='sources-wrap'>{pills}</div>", unsafe_allow_html=True)

def render_sources_table(lines: List[str]):
    if not lines:
        st.markdown("<p>No sources available.</p>", unsafe_allow_html=True)
        return

    table_html = '<style>.sources-table{width:auto;border-collapse:collapse;margin-top:0.5rem;}.sources-table th,.sources-table td{border:1px solid #e1efe4;padding:0.2rem;text-align:left;font-size:0.7rem;}.sources-table th{background:#f2f7f3;color:#0f3b1f;}.sources-table .link-cell{max-width:100px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;}.sources-table a{color:#0065BD;text-decoration:none;}.sources-table a:hover{text-decoration:underline;}</style><table class="sources-table"><thead><tr><th>Link</th><th>Location</th><th>ID</th></tr></thead><tbody>'

    for line in lines:
        parts = line.split(" — ", 2)
        if len(parts) < 3:
            continue
        link, location, cited_id = parts
        escaped_location = html.escape(location)
        escaped_id = html.escape(cited_id)
        full_link = globals().get('name_to_link', {}).get(link, link)
        if full_link.startswith('http'):
            escaped_link = html.escape(full_link)
            link_html = f'<a href="{escaped_link}" target="_blank" rel="noopener noreferrer">{escaped_link}</a>'
        else:
            link_html = html.escape(full_link)
        row_html = f'<tr><td class="link-cell">{link_html}</td><td>{escaped_location}</td><td>{escaped_id}</td></tr>'
        table_html += row_html

    table_html += '</tbody></table>'
    st.markdown(table_html, unsafe_allow_html=True)

def evaluate_top_sources(output: str, ranked: List[Tuple[Chunk, float]]) -> List[str]:
    if not os.environ.get("OPENAI_API_KEY", "").strip():
        return []

    context_blocks: List[str] = []
    for c, _ in ranked:
        context_blocks.append(f"ID: {c.id}, Source: {c.source_name}, Content: {c.text}")
    context = "\n\n".join(context_blocks)

    system_msg = "You are an AI evaluator. Given an output answer and chunks with sources, identify the top 3 source documents most crucial to generating the output. Return only the source names in order of relevance, separated by commas. No additional text."

    user_msg = f"Output: {output}\n\nChunks:\n{context}"

    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": user_msg},
    ]

    try:
        client = OpenAI()
        completion = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0,
            max_tokens=100,
        )
        response = completion.choices[0].message.content.strip()
        top = [s.strip() for s in response.split(',')][:3]
        return top
    except Exception:
        return []

def load_sources_reference_map(project_root: str) -> Dict[str, Dict[str, str]]:
    
    excel_path = os.path.join(project_root, 'reference', 'Sources - NEW 1-1.xlsx')
    if not os.path.exists(excel_path):
        return {}
    
    try:
        df = pd.read_excel(excel_path, engine='openpyxl')
    except Exception:
        return {}
    
    mapping: Dict[str, Dict[str, str]] = {}
    for _, row in df.iterrows():
        name_download = str(row.get('NAME DOWNLOAD', '')).lower().strip()
        if not name_download:
            continue
        title = str(row.get('Document Title', name_download)).strip()
        pr_ref = str(row.get('P&R Hub reference', '')).strip() if pd.notna(row.get('P&R Hub reference')) else ''
        mapping[name_download] = {
            'title': title,
            'ref': pr_ref
        }
    return mapping

def render_sources_leaderboard(ranked: List[Tuple[Chunk, float]], sources_ref_map: Dict[str, Dict[str, str]], top_sources: Optional[List[str]] = None):

    # Filter map to only include sources from Latest Updates
    latest_updates_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Latest Updates')
    latest_files = set(os.listdir(latest_updates_dir))
    filtered_map = {k: v for k, v in sources_ref_map.items() if k in latest_files or v['title'] in latest_files}

    if top_sources and len(top_sources) > 0:
        ordered = [(name, 1) for name in top_sources]
    else:
        if not ranked:
            return
        counts: Dict[str, int] = {}
        for c, _ in ranked:
            key = (c.source_name or '').strip()
            if not key:
                continue
            counts[key] = counts.get(key, 0) + 1
        if not counts:
            return

        ordered = sorted(counts.items(), key=lambda kv: kv[1], reverse=True)[:3]  # Limit to top 3

    rows: List[Tuple[str, str]] = []
    for name_download, _cnt in ordered:
        meta = filtered_map.get(name_download.lower()) if filtered_map else None
        title = meta['title'] if meta and meta.get('title') else name_download
        pr_ref = meta['ref'] if meta and meta.get('ref') else ''
        rows.append((title, pr_ref))

    if not rows:
        return

    # Build Markdown/HTML leaderboard with two columns and horizontal borders
    parts: List[str] = []
    parts.append('<style>.leaderboard{width:100%;border-collapse:collapse;margin:.5rem 0}.leaderboard thead th{font-size:.9rem;text-align:left;padding:.4rem 0;border-bottom:2px solid var(--border)}.leaderboard tbody tr{border-top:1px solid var(--border)}.leaderboard tbody td{padding:.5rem 0;font-size:.9rem}.leaderboard .title{font-weight:600}.leaderboard .ref{color:var(--muted)}</style>')
    parts.append('<table class="leaderboard"><thead><tr><th>Document Title</th><th>Reference</th></tr></thead><tbody>')
    for title, pr_ref in rows:
        safe_title = html.escape(title)
        safe_ref = html.escape(pr_ref) if pr_ref else '&#8212;'
        parts.append(f'<tr><td class="title">{safe_title}</td><td class="ref">{safe_ref}</td></tr>')
    parts.append('</tbody></table>')
    st.markdown("".join(parts), unsafe_allow_html=True)

def render_app() -> None:
    st.set_page_config(page_title="System Council Chatbot", page_icon="🌿", layout="centered")
    apply_cgiar_theme()

    DEFAULT_SYSTEM_PROMPT: str = (
        "You are a helpful assistant that answer user questions in a detailed and comprehensive way."
    )
    if 'system_prompt' not in st.session_state:
        st.session_state.system_prompt = DEFAULT_SYSTEM_PROMPT

    # Update the top image to a green nature scene, add frame
    st.markdown('<div style="height: 180px; border-radius: 12px; overflow: hidden; margin-bottom: 1rem; border: 4px solid #1F5A48;"><img src="https://images.unsplash.com/photo-1502082553048-f009c37129b9?auto=format&amp;fit=crop&amp;w=1000&amp;q=80" style="width:100%; height:100%; object-fit: cover;"></div>', unsafe_allow_html=True)

    # Add fixed top-left logo
    st.markdown("""
    <div style="position: fixed; top: 10px; left: 10px; z-index: 1000;">
        <img src="https://www.cgiar.org/wp/wp-content/uploads/2019/06/CGIAR_Logo.png" height="28">
    </div>
""", unsafe_allow_html=True)

    # Header with mini-logo
    st.markdown("""
        <div class="brand-hero">
            <h1 style="font-size: 2.375rem; text-align: center;">System Council Chatbot</h1>
            <p style="font-size: 1rem; text-align: justify;"><strong>PURPOSE: </strong>This ChatBot is designed to support stakeholders of the CGIAR System Council by providing fast, searchable access to governance materials on the System Council (mandates, meeting agendas and minutes, decisions and resolutions, policies, membership rules, and guidance). Its purpose is to help users find authoritative governance documents, interpret System Council decisions, locate relevant contacts and procedures, and streamline governance-related reporting and follow-up.</p>
            <p style="margin-bottom: 1rem;"></p>
            <p style="font-size: 1rem; text-align: justify;"><strong>ACKNOWLEDGMENT: </strong>This ChatBot uses Artificial Intelligence (AI) to understand questions and provide automated responses. While the models have been tested to ensure reliable and accurate information, AI-generated answers may occasionally contain errors or inaccuracies. If a response appears incorrect or unclear, please always refer to the official information available on the P&R Hub.</p>
        </div>
    """, unsafe_allow_html=True)

    # After the hero markdown, add a separator
    st.markdown("""<hr style="border-top: 9px solid #1F5A48; margin: 1rem 0;">""", unsafe_allow_html=True)

    # After separator, add back the loading logic
    force_reindex = False

    project_root = os.path.dirname(os.path.abspath(__file__))
    global name_to_link
    name_to_link = load_name_to_link(project_root)
    sources_ref_map = load_sources_reference_map(project_root)
    index_file = os.path.join(project_root, 'index.pkl')

    chunks = []
    embeddings = None
    regenerate = True

    if os.path.exists(index_file):
        with st.spinner("Loading from cache..."):
            try:
                with open(index_file, 'rb') as f:
                    data = pickle.load(f)
                if 'chunks' in data and 'embeddings' in data:
                    chunks = data['chunks']
                    embeddings = np.asarray(data['embeddings'])
                    if len(chunks) > 0 and embeddings.shape[0] == len(chunks):
                        regenerate = False
            except Exception:
                pass

    if regenerate:
        with st.spinner("Loading documents and building index..."):
            chunks = load_corpus(project_root)
            embeddings = build_embeddings_index(chunks)
            with open(index_file, 'wb') as f:
                pickle.dump({'chunks': chunks, 'embeddings': embeddings}, f)

    # Export chunks to Excel
    chunks_excel = os.path.join(project_root, 'chunks.xlsx')
    chunks_df = pd.DataFrame([c.__dict__ for c in chunks])
    chunks_df.to_excel(chunks_excel, index=False)

    num_chunks = len(chunks)

    if num_chunks == 0:
        st.warning(
            "No compatible documents (.pdf, .docx, .pptx) were found in the project. "
            "Add files to the existing folders and reload."
        )

    # Search area as a form
    with st.form("ask_form", clear_on_submit=False):
        query = st.text_input("Type your question:", value="", placeholder="e.g., What are the System Council’s recent decisions?", key="query_input")
        submitted = st.form_submit_button("🔎 Search", use_container_width=True)

    if submitted:
        loading_placeholder = st.empty()
        loading_placeholder.markdown("""
<style>
    @keyframes spinner-rotation {
        0% { transform: rotate(0deg); }
        100% { transform: rotate(360deg); }
    }
    @-webkit-keyframes spinner-rotation {
        0% { -webkit-transform: rotate(0deg); }
        100% { -webkit-transform: rotate(360deg); }
    }
    .loading-spinner {
        display: inline-block;
        width: 20px;
        height: 20px;
        border: 3px solid #1F5A48;
        border-radius: 50%;
        border-top-color: transparent;
        border-right-color: transparent;
        animation: spinner-rotation 1s linear infinite;
        -webkit-animation: spinner-rotation 1s linear infinite;
        vertical-align: middle;
        margin-right: 8px;
    }
</style>
<div style="display: flex; align-items: center; justify-content: center; margin: 1rem 0;">
    <span style="color: #1F5A48; font-weight: bold; font-size: 1.1rem;">
        <span class="loading-spinner"></span>In progress…
    </span>
</div>
""", unsafe_allow_html=True)

        top_k = 200
        model = get_embedding_model()
        ranked = rank_chunks(query, model, embeddings, chunks, top_k=top_k)

        # Try OpenAI (if API key present) per your rules
        ai_answer, openai_input = call_openai_generate(query, ranked, max_sentences=5, custom_system_msg=st.session_state.get('system_prompt'), name_to_link=name_to_link)

        top_sources = None

        # Answer card
        if ai_answer is None or not ai_answer.strip():
            answer, sources_all = compose_answer(query, ranked, name_to_link=name_to_link)

            # Keep answer compact
            unavailable = answer.startswith("I cannot find information in the provided chunks to answer this.") or \
                answer.startswith("Not found in the available information.")
            if not unavailable:
                sents = split_sentences(answer)
                if len(sents) > 6:
                    answer = " ".join(sents[:6]).strip()

            st.markdown(answer)

            used_for_leaderboard = ranked

        else:
            answer_part, used_ids = parse_used_sources(ai_answer)
            used_for_leaderboard = [(c, s) for c, s in ranked if c.id in used_ids]
            if not used_for_leaderboard:
                used_for_leaderboard = ranked  # fallback if no sources parsed
            st.markdown(answer_part)
            top_sources = evaluate_top_sources(answer_part, ranked)

        # Sources leaderboard (based on number of retrieved chunks per document)
        st.markdown("### TOP SOURCES")
        used_list = used_for_leaderboard if 'used_for_leaderboard' in locals() else ranked
        render_sources_leaderboard(used_list, sources_ref_map, top_sources=top_sources)

        # Removed old semantic sources table and 'Ver todas' per request

        # Logging interaction

        log_file = os.path.join(project_root, 'Logs', 'interaction_log.xlsx')

        ranked_chunk_details = ', '.join([f"{c.id}:{score:.3f}" for c, score in ranked])

        log_entry = {
            'timestamp': datetime.now().isoformat(),
            'query': query,
            'top_k': top_k,
            'num_ranked': len(ranked),
            'ai_used': bool(ai_answer and ai_answer.strip()),
            'system_prompt': st.session_state.system_prompt,
            'ranked_chunks': ranked_chunk_details,
            'answer': ai_answer if bool(ai_answer and ai_answer.strip()) else answer,
            'sources': ', '.join(format_sources_lines(ranked, max_items=3, name_to_link=name_to_link)),
            'retrieved_chunks_text': '\n\n---\n\n'.join(f"[ID: {c.id}]\n{c.text}" for c, _ in ranked),
            'openai_input': openai_input or ''
        }

        df_log = pd.DataFrame([log_entry])

        if os.path.exists(log_file):
            existing = pd.read_excel(log_file)
            df_log = pd.concat([existing, df_log], ignore_index=True)

        os.makedirs(os.path.dirname(log_file), exist_ok=True)
        df_log.to_excel(log_file, index=False)

        loading_placeholder.empty()

    st.markdown('</div>', unsafe_allow_html=True)

    # Footer
    st.markdown(
        """<div class='app-footer'>If we were unable to answer your question or you require further information, please contact us at: <a href='mailto:@cgiar.org'>@cgiar.org</a></div>""",
        unsafe_allow_html=True
    )


if __name__ == "__main__":
    render_app()